"""
Generate PowerPoint Presentation 2: Implementation & Intermediate Results
For 3D Point Cloud Reconstruction Project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(25, 51, 102)  # Dark blue
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.word_wrap = True
    p = subtitle_frame.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(100, 181, 246)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    # Add blue underline
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
        p.level = 0
    
    return slide

def add_two_column_slide(prs, title, left_content, right_content):
    """Add a two-column content slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    # Blue underline
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_content):
        if i == 0:
            p = left_frame.paragraphs[0]
        else:
            p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_content):
        if i == 0:
            p = right_frame.paragraphs[0]
        else:
            p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    return slide

def create_presentation():
    """Create the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # SLIDE 1: Title Slide
    add_title_slide(prs, 
                    "3D Point Cloud Reconstruction",
                    "Progress Update: Implementation & Intermediate Results")
    
    # SLIDE 2: Where We Left Off
    add_content_slide(prs,
                      "Where We Left Off",
                      [
                          "• Problem: Automatic 3D reconstruction of industrial parts from monocular video requires robust feature matching and multi-view geometry",
                          "",
                          "• Dataset: 78 configurations across 3 complexity levels (single, multiple, stacked parts) from industrial scanning setup",
                          "",
                          "• Planned Approach: SIFT → Feature Matching → Essential Matrix → Triangulation → Point Cloud refinement"
                      ])
    
    # SLIDE 3: Pipeline Implementation Status
    add_content_slide(prs,
                      "Pipeline Implementation Status",
                      [
                          "✅ Frame Extraction & Undistortion [100%]",
                          "✅ SIFT Feature Detection [100%]",
                          "✅ Feature Matching (Lowe's ratio test) [100%]",
                          "✅ Essential Matrix Computation (RANSAC) [100%]",
                          "✅ Camera Pose Recovery [100%]",
                          "✅ Triangulation & Point Cloud Creation [100%]",
                          "✅ Basic Filtering (Outlier Removal) [100%]",
                          "",
                          "🔄 Advanced Filtering & Optimization [60%]",
                          "🔄 Quantitative Evaluation (all 78 configs) [45%]",
                          "⏳ Final Refinement & Results Analysis [0%]"
                      ])
    
    # SLIDE 4: Algorithm Improvements
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Algorithm Improvements Since Presentation 1"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    improvements = [
        ("Feature Detector", "SIFT (slower)", "SuperPoint + LightGlue OR LoFTR"),
        ("Matrix Estimation", "Fundamental (7 DOF)", "Essential Matrix (5 DOF)"),
        ("Motion Constraint", "General motion", "Homography (4 DOF) - fixed camera"),
        ("Reconstruction Scale", "Unknown scale", "Metric - direct from calibrated camera")
    ]
    
    for i, (aspect, before, after) in enumerate(improvements):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = f"{aspect}: {before} → {after}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    return slide, prs

# SLIDE 5: Advanced Triangulation
def add_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Advanced Triangulation Techniques"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    content = [
        "Approach A: Monocular Depth Priors",
        "  • Tools: Depth Anything v2 or UniDepth",
        "  • Status: ✅ Implemented as optional pre-processor",
        "",
        "Approach B: Multi-View Stereo (MVS) Networks",
        "  • Tools: CasMVSNet or IterMVS",
        "  • Benefit: 10-50× denser than classical triangulation",
        "  • Status: 🔄 Integration in progress",
        "",
        "Approach C: Direct Regression (DuSt3R / MASt3R)",
        "  • Transformer-based: Image pair → 3D point map",
        "  • Advantage: Bypasses feature matching",
        "  • Status: ⏳ Planned for final refinement"
    ]
    
    for i, item in enumerate(content):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    
    return slide

# SLIDE 6: Bugs & Challenges
def add_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Bugs & Challenges Faced"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    challenges = [
        "Challenge 1: Multi-Object Degradation ⚠️",
        "  Problem: Quality drops with 2+ parts",
        "  Solution: Stricter ratio test, per-object pose estimation",
        "",
        "Challenge 2: Fingerprint Noise (Clay Capacitor) ✅",
        "  Problem: Fingerprints picked up as 3D points",
        "  Solution: Statistical filtering + lighting adjustment",
        "  Status: 70% noise removed",
        "",
        "Challenge 3: Low Point Count (Some Configs) ✅",
        "  Problem: <50K points on certain parts",
        "  Solution: Increased frame sampling, better initialization",
        "  Status: Now generating 200K-470K points"
    ]
    
    for i, item in enumerate(challenges):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    
    return slide

# SLIDE 7: Demo
def add_slide_7(prs):
    add_content_slide(prs,
                      "Pipeline in Action",
                      [
                          "Input: 500+ frames from 4K video (G-LS-I-LO-33)",
                          "",
                          "Processing Steps:",
                          "  1. Feature Detection: ~2000 features per frame",
                          "  2. Feature Matching: 110K+ keypoint pairs identified",
                          "  3. Camera Geometry: Essential matrix via RANSAC",
                          "  4. Triangulation: 2D matches → 3D points",
                          "",
                          "Output: 208K point cloud rendered in CloudCompare",
                          "",
                          "Processing Time: 2-3 minutes end-to-end"
                      ])

# SLIDE 8: Intermediate Results
def add_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Early Point Cloud Results"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    content = [
        "Example Point Clouds (Level 1 - Single Objects):",
        "",
        "Config              | Points     | Matches  | Status",
        "G-LS-I-LO-33        | 208,046    | 110,932  | ✅ Excellent",
        "G-MS-I-LO-1         | 233,822    | 68,440   | ✅ Good",
        "G-TS-P-HI-3         | 473,042    | 180,642  | ✅ Very Good",
        "",
        "Early Quantitative Metrics:",
        "  • Completeness: 65-78% vs ground truth",
        "  • Accuracy: 3-7mm average error",
        "  • Noise Level: 12-18% false positives",
        "",
        "⚠️ Full 78-config evaluation in progress (ETA: 2-3 days)"
    ]
    
    for i, item in enumerate(content):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    
    return slide

# SLIDE 9: What's Still Pending
def add_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "Remaining Work for Presentation 3"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = RGBColor(25, 51, 102)
    
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(9), Inches(0))
    line.line.color.rgb = RGBColor(100, 181, 246)
    line.line.width = Pt(3)
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(8.8), Inches(5.5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    remaining = [
        "🔴 CRITICAL (Blocking Full Evaluation):",
        "  1. Complete all 78 configurations (ETA: 3-4 days)",
        "  2. Resolve multi-object degradation",
        "  3. Fix remaining 30% fingerprint noise",
        "",
        "🟡 IMPORTANT (Improve Results Quality):",
        "  1. Integrate MVS networks for denser clouds",
        "  2. Implement depth prior alignment",
        "  3. Optimize adaptive voxel downsampling",
        "  4. Verify camera intrinsics calibration",
        "",
        "Timeline for Presentation 3: ~4 weeks",
        "Expected: Complete, optimized pipeline + comprehensive evaluation"
    ]
    
    for i, item in enumerate(remaining):
        p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(50, 50, 50)
        p.space_before = Pt(2)
        p.space_after = Pt(2)
    
    return slide

if __name__ == "__main__":
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Add all slides
    add_title_slide(prs, 
                    "3D Point Cloud Reconstruction",
                    "Progress Update: Implementation & Intermediate Results")
    
    add_content_slide(prs,
                      "Where We Left Off",
                      [
                          "• Problem: Automatic 3D reconstruction of industrial parts from monocular video requires robust feature matching and multi-view geometry",
                          "",
                          "• Dataset: 78 configurations across 3 complexity levels (single, multiple, stacked parts) from industrial scanning setup",
                          "",
                          "• Planned Approach: SIFT → Feature Matching → Essential Matrix → Triangulation → Point Cloud refinement"
                      ])
    
    add_content_slide(prs,
                      "Pipeline Implementation Status",
                      [
                          "✅ Frame Extraction & Undistortion [100%]",
                          "✅ SIFT Feature Detection [100%]",
                          "✅ Feature Matching (Lowe's ratio test) [100%]",
                          "✅ Essential Matrix Computation (RANSAC) [100%]",
                          "✅ Camera Pose Recovery [100%]",
                          "✅ Triangulation & Point Cloud Creation [100%]",
                          "✅ Basic Filtering (Outlier Removal) [100%]",
                          "",
                          "🔄 Advanced Filtering & Optimization [60%]",
                          "🔄 Quantitative Evaluation (all 78 configs) [45%]",
                          "⏳ Final Refinement & Results Analysis [0%]"
                      ])
    
    slide4, prs = create_presentation()
    add_slide_5(prs)
    add_slide_6(prs)
    add_slide_7(prs)
    add_slide_8(prs)
    add_slide_9(prs)
    
    # Save presentation
    output_file = "Presentation_2_Implementation_and_Intermediate_Results.pptx"
    prs.save(output_file)
    print(f"✅ Presentation created successfully: {output_file}")
    print(f"📊 Total slides: {len(prs.slides)}")
